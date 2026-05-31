#!/usr/bin/env python3
"""
Build Private Marketing Agency Blueprint — PowerPoint presentation (Thai)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Dimensions: 16:9 widescreen ──────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Brand colours ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x17)   # deep dark
SURFACE   = RGBColor(0x1E, 0x25, 0x33)   # card bg
SURFACE2  = RGBColor(0x16, 0x1B, 0x22)   # alt card bg
BORDER    = RGBColor(0x2D, 0x37, 0x48)   # subtle line
VIOLET    = RGBColor(0x7C, 0x3A, 0xED)   # primary accent
VIOLET_L  = RGBColor(0x9F, 0x5F, 0xFF)   # light violet
CYAN      = RGBColor(0x06, 0xB6, 0xD4)   # secondary accent
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # success
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # warning
RED       = RGBColor(0xEF, 0x44, 0x44)   # danger
TEXT      = RGBColor(0xF1, 0xF5, 0xF9)   # primary text
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # secondary text
CODE      = RGBColor(0xA5, 0xF3, 0xFC)   # monospace highlight
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)


# ── Helper: add background rectangle ─────────────────────────────────────────
def fill_bg(slide, color=BG):
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.zorder = 0
    return bg


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_label(slide, text, left, top, width=Inches(3), size=10, color=MUTED, bold=False, align=PP_ALIGN.LEFT):
    return add_text(slide, text, left, top, width, Inches(0.4), size, bold, color, align)


# Left margin / padding constants
LM = Inches(0.6)
TM = Inches(0.45)
CW = SLIDE_W - Inches(1.2)   # content width


def heading(slide, text, top=TM, size=32, color=TEXT):
    add_text(slide, text, LM, top, CW, Inches(1.0), size, True, color)


def subheading(slide, text, top=Inches(1.2), size=14, color=MUTED):
    add_text(slide, text, LM, top, CW, Inches(0.5), size, False, color)


def tagline_bar(slide, top=Inches(1.45), color=VIOLET):
    bar = add_rect(slide, LM, top, Inches(1.0), Inches(0.06), color)
    return bar


def section_pill(slide, text, top=TM - Inches(0.1)):
    """Small pill label at top left"""
    box = add_rect(slide, LM, top, Inches(2.2), Inches(0.35), SURFACE, VIOLET)
    add_text(slide, text, LM + Inches(0.1), top + Inches(0.05),
             Inches(2.0), Inches(0.25), 9, True, VIOLET_L, PP_ALIGN.LEFT)


def bullet_item(slide, text, left, top, width, color=TEXT,
                prefix="—", prefix_color=CYAN, size=13):
    """Single bullet line with prefix dash"""
    # prefix
    add_text(slide, prefix, left, top, Inches(0.3), Inches(0.4),
             size, True, prefix_color)
    # text
    add_text(slide, text, left + Inches(0.32), top, width - Inches(0.35),
             Inches(0.45), size, False, color, PP_ALIGN.LEFT)


def checklist_item(slide, text, left, top, width, check=True, size=13):
    """Checklist item with ✓ or ✗"""
    sym = "✓" if check else "✗"
    col = GREEN if check else RED
    add_text(slide, sym, left, top, Inches(0.32), Inches(0.4), size, True, col)
    add_text(slide, text, left + Inches(0.35), top, width - Inches(0.38),
             Inches(0.45), size, False, TEXT, PP_ALIGN.LEFT)


def card(slide, left, top, width, height, border_color=BORDER, fill=SURFACE):
    r = add_rect(slide, left, top, width, height, fill, border_color)
    return r


# ─────────────────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


def new_slide():
    s = prs.slides.add_slide(blank_layout)
    fill_bg(s)
    return s


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Hero
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
# Glow rectangle (decorative)
g = add_rect(s, Inches(8.5), Inches(-1), Inches(6), Inches(6),
             RGBColor(0x1A, 0x0A, 0x2E), None)

add_text(s, "Claude Code Course  ·  8 บทเรียน  ·  ~2.5 ชั่วโมง",
         LM, TM, CW, Inches(0.4), 11, False, VIOLET_L)

add_text(s, "Private Marketing\nAgency Blueprint",
         LM, Inches(1.1), Inches(8), Inches(2.2), 48, True, TEXT)

tagline_bar(s, Inches(3.5))

add_text(s, "สร้างระบบ AI Marketing ของคุณเองในวันเดียว",
         LM, Inches(3.7), Inches(8), Inches(0.6), 18, False, MUTED)

# Pills row
pills = [
    ("ไม่ต้องเขียนโค้ด", VIOLET_L),
    ("ไม่ต้องจ้างทีม", GREEN),
    ("ผลลัพธ์จริง ใช้งานได้ทันที", CYAN),
]
px = LM
for txt, col in pills:
    w = Inches(2.4)
    add_rect(s, px, Inches(4.4), w, Inches(0.38), SURFACE, col)
    add_text(s, txt, px + Inches(0.12), Inches(4.45), w - Inches(0.12),
             Inches(0.3), 11, True, col, PP_ALIGN.LEFT)
    px += w + Inches(0.15)

add_text(s, "Claude Code · OpenRouter · NocoDB · HeyGen · Vercel",
         Inches(7.5), Inches(6.8), Inches(5.5), Inches(0.4),
         10, False, MUTED, PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Pain Points
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "คุณใช้เวลากี่ชั่วโมงต่อสัปดาห์กับการตลาด?")
tagline_bar(s)

pain = [
    ("เขียน caption ทุกวัน", "ซ้ำแล้วซ้ำเล่า ไม่รู้ว่าหัวข้อไหนโดน ต้องนึกเองทุกครั้ง"),
    ("ออกแบบกราฟิก", "ทั้งที่ไม่ได้เป็นดีไซเนอร์ ใช้ Canva วนไปวนมา"),
    ("วางแผนคอนเทนต์", "ไม่มีระบบ วางแผนแบบ ad hoc ขาดความต่อเนื่อง"),
    ("ทำทุกอย่างคนเดียว", "วิดีโอ แลนดิ้งเพจ เว็บไซต์ งานการตลาดไม่มีจบ"),
]

col_w  = Inches(5.9)
col_h  = Inches(1.5)
gap    = Inches(0.2)
row_y  = [Inches(1.7), Inches(3.3)]
col_x  = [LM, LM + col_w + gap]

for i, (title, desc) in enumerate(pain):
    x = col_x[i % 2]
    y = row_y[i // 2]
    card(s, x, y, col_w, col_h, VIOLET)
    add_text(s, title, x + Inches(0.2), y + Inches(0.15),
             col_w - Inches(0.3), Inches(0.45), 15, True, TEXT)
    add_text(s, desc, x + Inches(0.2), y + Inches(0.65),
             col_w - Inches(0.3), Inches(0.7), 12, False, MUTED)

# Callout
add_rect(s, LM, Inches(5.0), CW, Inches(0.55), RGBColor(0x1A, 0x0A, 0x2E), VIOLET)
add_text(s, "นี่คือปัญหาของเจ้าของ SME และ solopreneur ทุกคน ไม่ใช่แค่คุณ",
         LM + Inches(0.2), Inches(5.08), CW - Inches(0.3), Inches(0.4),
         13, False, VIOLET_L)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tools don't connect
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "เครื่องมือ AI มีเยอะ  แต่ไม่คุยกัน")
tagline_bar(s)

tools = [
    ("ChatGPT",    Inches(1.2), Inches(1.8)),
    ("Canva",      Inches(4.5), Inches(2.1)),
    ("CapCut",     Inches(8.0), Inches(1.7)),
    ("Notion",     Inches(2.5), Inches(3.2)),
    ("Buffer",     Inches(6.2), Inches(3.0)),
    ("Midjourney", Inches(9.5), Inches(2.5)),
    ("Google Docs",Inches(1.0), Inches(4.2)),
    ("HeyGen",     Inches(5.3), Inches(4.5)),
    ("Airtable",   Inches(9.2), Inches(4.0)),
]
for name, tx, ty in tools:
    w = Inches(1.9)
    card(s, tx, ty, w, Inches(0.45), BORDER)
    add_text(s, name, tx + Inches(0.12), ty + Inches(0.08), w - Inches(0.15),
             Inches(0.32), 12, False, MUTED)

# Big ? in the middle
add_text(s, "?", Inches(5.8), Inches(2.8), Inches(1.5), Inches(1.0),
         60, True, AMBER, PP_ALIGN.CENTER)

# Callout bottom
add_rect(s, LM, Inches(5.5), CW, Inches(0.6), RGBColor(0x1A, 0x10, 0x00), AMBER)
add_text(s, "แอปนึงสำหรับเขียน  ·  แอปนึงสำหรับรูป  ·  แอปนึงสำหรับวิดีโอ — ไม่มีอะไรเชื่อมกัน",
         LM + Inches(0.2), Inches(5.6), CW - Inches(0.3), Inches(0.42),
         13, False, AMBER, PP_ALIGN.LEFT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What to automate
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "ไม่ใช่ทุกงานที่ควร Automate")
tagline_bar(s)

# Left column: hard
add_rect(s, LM, Inches(1.7), Inches(5.9), Inches(3.1),
         RGBColor(0x2A, 0x0E, 0x0E), RGBColor(0xEF, 0x44, 0x44))
add_text(s, "Automate ยาก", LM + Inches(0.2), Inches(1.85),
         Inches(5.5), Inches(0.4), 14, True, RGBColor(0xFC, 0xA5, 0xA5))
hard = ["ต้องตัดสินใจเชิงกลยุทธ์", "เปลี่ยนแปลงทุกครั้ง",
        "ต้องใช้ความสัมพันธ์กับลูกค้า", "ต้องใช้ประสบการณ์เฉพาะทาง"]
for i, t in enumerate(hard):
    bullet_item(s, t, LM + Inches(0.1), Inches(2.3) + Inches(0.5) * i,
                Inches(5.6), MUTED, "—", RGBColor(0xEF, 0x44, 0x44), 13)

# Right column: easy
rx = LM + Inches(6.1)
add_rect(s, rx, Inches(1.7), Inches(5.9), Inches(3.1),
         RGBColor(0x06, 0x20, 0x17), GREEN)
add_text(s, "Automate ได้ดีมาก", rx + Inches(0.2), Inches(1.85),
         Inches(5.5), Inches(0.4), 14, True, GREEN)
easy = ["ทำซ้ำเดิมทุกครั้ง", "มีโครงสร้างชัดเจน",
        "Production มากกว่า Judgement", "ใช้ข้อมูลตายตัวเป็น input"]
for i, t in enumerate(easy):
    bullet_item(s, t, rx + Inches(0.1), Inches(2.3) + Inches(0.5) * i,
                Inches(5.6), TEXT, "—", GREEN, 13)

# Callout
add_rect(s, LM, Inches(5.1), CW, Inches(0.55), RGBColor(0x12, 0x0A, 0x28), VIOLET)
add_text(s, "การตลาดคอนเทนต์ = perfect case สำหรับ automation",
         LM + Inches(0.2), Inches(5.18), CW - Inches(0.3), Inches(0.38),
         13, True, VIOLET_L, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Vision
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_text(s, "จินตนาการว่า...",
         LM, Inches(1.2), CW, Inches(0.4), 12, False, MUTED)

add_text(s, "คุณพิมพ์คำสั่งเดียว",
         LM, Inches(1.8), CW, Inches(0.65), 28, True, TEXT)

add_text(s, "ระบบ AI วิจัย วางแผน เขียน ออกแบบ\nทำวิดีโอ และ deploy เว็บ — โดยอัตโนมัติ",
         LM, Inches(2.55), CW, Inches(1.4), 28, True, VIOLET_L)

tagline_bar(s, Inches(4.15), CYAN)

add_text(s, "นั่นคือสิ่งที่คุณจะสร้างในคอร์สนี้ — โดยใช้แบรนด์ของคุณเอง สำหรับธุรกิจของคุณเอง",
         LM, Inches(4.4), CW, Inches(0.5), 14, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Section Header: Solution
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, VIOLET, None)
add_text(s, "ส่วนที่ 01", LM, Inches(2.0), CW, Inches(0.45), 14, True, VIOLET_L)
add_text(s, "ระบบที่คุณจะสร้าง", LM, Inches(2.55), CW, Inches(1.2), 40, True, TEXT)
tagline_bar(s, Inches(3.9))
add_text(s, "Private Marketing Agency ใน Claude Code",
         LM, Inches(4.1), CW, Inches(0.5), 16, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — What is Claude Code
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Claude Code คืออะไร")
tagline_bar(s)
subheading(s, "ผู้ช่วย AI ที่ทำงานได้จริงในเครื่องของคุณ — ไม่ใช่แค่ chatbot")

points = [
    ("อ่านและสร้างไฟล์ได้จริง",
     "รู้ว่าโปรเจกต์คุณมีไฟล์อะไร ทำงานกับโฟลเดอร์ได้โดยตรง"),
    ("เชื่อมต่อกับระบบภายนอกได้",
     "ต่อกับ database, video platform, web hosting ผ่านสิ่งที่เรียกว่า MCP"),
    ("ไม่ต้องเขียนโค้ดเอง",
     "คุณอธิบายสิ่งที่ต้องการ Claude เป็นคนสร้างทุกอย่าง"),
]
for i, (title, desc) in enumerate(points):
    y = Inches(1.85) + Inches(1.5) * i
    card(s, LM, y, CW, Inches(1.3), VIOLET)
    add_text(s, title, LM + Inches(0.2), y + Inches(0.18),
             CW - Inches(0.3), Inches(0.4), 15, True, TEXT)
    add_text(s, desc, LM + Inches(0.2), y + Inches(0.62),
             CW - Inches(0.3), Inches(0.55), 12, False, MUTED)

add_rect(s, LM, Inches(6.45), CW, Inches(0.5),
         RGBColor(0x06, 0x1A, 0x20), CYAN)
add_text(s, "คุณเปิด folder ใน Claude Code แล้วพิมพ์คำสั่ง — Claude ทำงานให้",
         LM + Inches(0.2), Inches(6.54), CW - Inches(0.3), Inches(0.35),
         13, False, CYAN, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Pipeline overview
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Private Marketing Agency = 8 Agents ทำงานต่อกัน")
tagline_bar(s)

steps = ["วิจัย", "วางแผน", "เขียน", "ออกแบบ", "วิดีโอ", "เว็บ", "เผยแพร่"]
n = len(steps)
box_w = Inches(1.55)
gap_w = Inches(0.1)
total = box_w * n + gap_w * (n - 1) + Inches(0.28) * (n - 1)  # arrows
# fit all in slide width
start_x = LM
arrow_w = Inches(0.28)
unit = (CW - box_w * n - arrow_w * (n-1)) / max(n - 1, 1)

x = start_x
y_box = Inches(2.2)
box_h = Inches(1.6)

for i, label in enumerate(steps):
    card(s, x, y_box, box_w, box_h, VIOLET, SURFACE)
    add_text(s, label, x, y_box + Inches(0.55), box_w,
             Inches(0.6), 16, True, CYAN, PP_ALIGN.CENTER)
    x += box_w
    if i < n - 1:
        add_text(s, "→", x, y_box + Inches(0.55), arrow_w,
                 Inches(0.6), 22, True, VIOLET, PP_ALIGN.CENTER)
        x += arrow_w

add_text(s, "ทั้งหมดทำงานอยู่ใน Claude Code ของคุณ ต่อกันเป็น pipeline เดียว",
         LM, Inches(4.1), CW, Inches(0.45), 14, False, MUTED)
add_text(s, "แต่ละขั้นตอนรู้จัก brand ของคุณ และส่งต่อ output ให้ขั้นตอนถัดไป",
         LM, Inches(4.6), CW, Inches(0.45), 14, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "เครื่องมือที่ใช้ในคอร์สนี้")
tagline_bar(s)

# Left col: Core
add_text(s, "CORE", LM, Inches(1.7), Inches(5.9), Inches(0.4),
         10, True, MUTED)
core = [
    ("Claude Code", "AI conductor — รู้จักโปรเจกต์ รันคำสั่ง เชื่อมทุกอย่าง", VIOLET),
    ("OpenRouter",  "Gateway สำหรับ AI models — สร้างรูป, ค้นหา model ใหม่ๆ", CYAN),
]
cy = Inches(2.05)
for name, desc, col in core:
    card(s, LM, cy, Inches(5.9), Inches(1.0), col)
    add_text(s, name, LM + Inches(0.2), cy + Inches(0.1),
             Inches(5.5), Inches(0.4), 15, True, TEXT)
    add_text(s, desc, LM + Inches(0.2), cy + Inches(0.5),
             Inches(5.5), Inches(0.4), 12, False, MUTED)
    cy += Inches(1.1)

# Right col: MCPs
rx = LM + Inches(6.2)
add_text(s, "INTEGRATIONS (MCP)", rx, Inches(1.7), Inches(5.9), Inches(0.4),
         10, True, MUTED)
mcps = [
    ("NocoDB",  "Database สำหรับ content calendar"),
    ("HeyGen",  "AI avatar video — script → วิดีโอ"),
    ("Vercel",  "Web hosting — deploy ใน 1 คำสั่ง"),
]
cy = Inches(2.05)
for name, desc in mcps:
    card(s, rx, cy, Inches(5.9), Inches(0.88), BORDER)
    add_text(s, name, rx + Inches(0.2), cy + Inches(0.08),
             Inches(5.5), Inches(0.38), 14, True, TEXT)
    add_text(s, desc, rx + Inches(0.2), cy + Inches(0.46),
             Inches(5.5), Inches(0.35), 12, False, MUTED)
    cy += Inches(0.98)

add_text(s, "ทุกอย่างตั้งค่าระหว่างเรียน ไม่ต้องเตรียมล่วงหน้า",
         LM, Inches(5.5), CW, Inches(0.4), 12, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Deliverables
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "ผลลัพธ์ที่จับต้องได้  เมื่อจบ 2.5 ชั่วโมง")
tagline_bar(s)

items_l = [
    "โฟลเดอร์โปรเจกต์ที่ Claude รู้จักแบรนด์คุณทุก session",
    "Brand Guidelines ที่ agents ทุกตัวอ่าน",
    "Content Calendar ใน NocoDB (20–30 รายการ)",
    "บทความ, caption, video script เสียงเหมือนคุณ",
]
items_r = [
    "ภาพกราฟิกที่ตรงกับสไตล์แบรนด์",
    "วิดีโอ AI Avatar ที่พูดสคริปต์คุณ",
    "Landing page ที่ deploy ขึ้น internet จริง",
    "คำสั่ง /content-engine รันทั้ง pipeline ด้วยคีย์เดียว",
]
for i, txt in enumerate(items_l):
    checklist_item(s, txt, LM, Inches(1.85) + Inches(0.88) * i, Inches(5.9))
for i, txt in enumerate(items_r):
    checklist_item(s, txt, LM + Inches(6.2), Inches(1.85) + Inches(0.88) * i, Inches(5.9))

add_rect(s, LM, Inches(5.65), CW, Inches(0.55), RGBColor(0x06, 0x20, 0x17), GREEN)
add_text(s, "ทุกชิ้นเป็นของแบรนด์คุณ ใช้งานได้จริงทันทีหลังจบคอร์ส",
         LM + Inches(0.2), Inches(5.73), CW - Inches(0.3), Inches(0.38),
         13, False, GREEN, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Section Header: Lessons
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, VIOLET, None)
add_text(s, "ส่วนที่ 02", LM, Inches(2.0), CW, Inches(0.45), 14, True, VIOLET_L)
add_text(s, "8 บทเรียน\nสร้างทีละชิ้น", LM, Inches(2.55), CW, Inches(1.8), 40, True, TEXT)
tagline_bar(s, Inches(4.55))
add_text(s, "แต่ละบทเรียนผลิต output จริง ที่บทถัดไปใช้ต่อ",
         LM, Inches(4.75), CW, Inches(0.5), 16, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Timeline overview
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "ภาพรวม 8 บทเรียน")
tagline_bar(s)

lessons_tl = [
    ("01", "Brand\nContext",  "20 min"),
    ("02", "Researcher\nAgent", "20 min"),
    ("03", "Content\nCalendar", "20 min"),
    ("04", "Content\nWriter", "15 min"),
    ("05", "Content\nDesigner", "20 min"),
    ("06", "Video\nMaker", "20 min"),
    ("07", "Landing\nPage", "20 min"),
    ("08", "Publisher\n+ Wrap", "15 min"),
]
n = len(lessons_tl)
circle_d = Inches(0.75)
total_w  = CW
spacing  = total_w / n
cx_start = LM + spacing / 2

# Connector line
line_y = Inches(2.45)
add_rect(s, LM, line_y, CW, Inches(0.04), VIOLET, None)

for i, (num, label, time_) in enumerate(lessons_tl):
    cx = cx_start + spacing * i
    # circle
    circ = add_rect(s, cx - circle_d/2, line_y - circle_d/2,
                    circle_d, circle_d, VIOLET, None)
    add_text(s, num, cx - circle_d/2, line_y - circle_d/2 + Inches(0.12),
             circle_d, circle_d - Inches(0.12), 14, True, WHITE, PP_ALIGN.CENTER)
    # label below
    add_text(s, label,
             cx - Inches(0.8), line_y + circle_d/2 + Inches(0.08),
             Inches(1.6), Inches(0.75), 11, False, TEXT, PP_ALIGN.CENTER)
    # time above
    add_text(s, time_,
             cx - Inches(0.7), line_y - circle_d/2 - Inches(0.38),
             Inches(1.4), Inches(0.32), 10, False, MUTED, PP_ALIGN.CENTER)

# Total pill
add_rect(s, SLIDE_W - Inches(2.4), Inches(6.4), Inches(1.95), Inches(0.42),
         SURFACE, VIOLET)
add_text(s, "รวม ~2.5 ชั่วโมง",
         SLIDE_W - Inches(2.35), Inches(6.48), Inches(1.85), Inches(0.3),
         11, True, VIOLET_L, PP_ALIGN.LEFT)


# ═════════════════════════════════════════════════════════════════════════════
# Helper: lesson card pair (2 cards side by side)
# ═════════════════════════════════════════════════════════════════════════════
def lesson_card_pair(prs, num_a, title_a, items_a, out_a,
                     num_b, title_b, items_b, out_b, heading_text, callout=None):
    s = new_slide()
    heading(s, heading_text)
    tagline_bar(s)

    cw2 = Inches(5.9)
    for xi, (num, title, items, out) in enumerate([
        (num_a, title_a, items_a, out_a),
        (num_b, title_b, items_b, out_b)
    ]):
        x = LM if xi == 0 else LM + Inches(6.2)
        y_top = Inches(1.7)
        card(s, x, y_top, cw2, Inches(5.1), BORDER)
        # Badge
        add_text(s, num, x + Inches(0.15), y_top + Inches(0.12),
                 Inches(1.0), Inches(0.55), 28, True, VIOLET_L)
        # Title
        add_text(s, title, x + Inches(0.15), y_top + Inches(0.72),
                 cw2 - Inches(0.25), Inches(0.42), 15, True, TEXT)
        # Divider
        add_rect(s, x + Inches(0.1), y_top + Inches(1.2), cw2 - Inches(0.2),
                 Inches(0.02), BORDER, None)
        # Bullet items
        for j, item in enumerate(items):
            add_text(s, "›  " + item,
                     x + Inches(0.15), y_top + Inches(1.32) + Inches(0.55) * j,
                     cw2 - Inches(0.25), Inches(0.5), 12, False, MUTED)
        # Output footer
        add_rect(s, x + Inches(0.1), y_top + Inches(4.35), cw2 - Inches(0.2),
                 Inches(0.02), BORDER, None)
        add_text(s, "Output: " + out,
                 x + Inches(0.15), y_top + Inches(4.45), cw2 - Inches(0.25),
                 Inches(0.5), 10, False, CYAN)

    if callout:
        add_rect(s, LM, Inches(6.95), CW, Inches(0.42),
                 RGBColor(0x06, 0x1A, 0x20), CYAN)
        add_text(s, callout, LM + Inches(0.2), Inches(7.02),
                 CW - Inches(0.3), Inches(0.3), 12, False, CYAN)
    return s


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Lessons 1-2
# ═════════════════════════════════════════════════════════════════════════════
lesson_card_pair(
    prs,
    "01", "Brand Context",
    ["สร้าง CLAUDE.md — briefing ถาวรสำหรับทุก session",
     "ตั้งค่า OpenRouter plugin",
     "สร้าง brand-guidelines.md ด้วย 6 คำถามกุญแจ",
     "สร้างคำสั่ง /brand"],
    "CLAUDE.md, brand/brand-guidelines.md, /brand command",
    "02", "Researcher Agent",
    ["สร้าง SKILL.md สำหรับ content-researcher",
     "ค้นหา trending topics ล่าสุด (30 วัน)",
     "วิเคราะห์ competitor content, audience questions",
     "สร้างคำสั่ง /research"],
    ".claude/skills/content-researcher/SKILL.md, research/ folder",
    "Lessons 1–2",
    "ทุกบทเรียนผลิตไฟล์จริง — ไม่ใช่แค่ความรู้ที่หายไปหลังปิดหน้าต่าง"
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Lessons 3-4
# ═════════════════════════════════════════════════════════════════════════════
lesson_card_pair(
    prs,
    "03", "Content Calendar",
    ["ตั้งค่า NocoDB MCP — ต่อ Claude กับ database",
     "สร้าง Content Calendar table (7 fields)",
     "Generate 20–30 รายการสำหรับเดือนถัดไป",
     "รายการกระจายหลาย platform และ format"],
    "NocoDB table พร้อม 20+ entries ครอบทุก pillar",
    "04", "Content Writer",
    ["สร้าง SKILL.md สำหรับ content-writer",
     "เขียน social post ในเสียงของแบรนด์",
     "เขียน article 800–1,500 คำ",
     "เขียน video script ที่พูดได้จริง"],
    "3+ ไฟล์ใน content/ folder, NocoDB status = Drafted",
    "Lessons 3–4",
    "Key: เนื้อหาต้องฟังดูเหมือนคุณ ไม่ใช่ generic AI marketing copy"
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Lessons 5-6
# ═════════════════════════════════════════════════════════════════════════════
lesson_card_pair(
    prs,
    "05", "Content Designer",
    ["สร้าง SKILL.md สำหรับ content-designer",
     "Generate รูปที่ตรงสีและสไตล์แบรนด์",
     "ใช้ OpenRouter image generation",
     "Aspect ratio ถูก platform: 1:1, 9:16, 16:9"],
    "2+ รูปใน designs/ folder, NocoDB status = Designed",
    "06", "Video Maker",
    ["ตั้งค่า HeyGen MCP — ต่อกับ AI avatar platform",
     "เลือก avatar และ voice ที่เหมาะกับแบรนด์",
     "สร้าง SKILL.md สำหรับ video-maker",
     "Generate วิดีโอ (render 2–5 นาที)"],
    "วิดีโอ AI avatar พูด script ของคุณ 1 ชิ้น",
    "Lessons 5–6"
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Lessons 7-8
# ═════════════════════════════════════════════════════════════════════════════
lesson_card_pair(
    prs,
    "07", "Landing Page",
    ["ตั้งค่า Vercel MCP — ต่อกับ web hosting",
     "Claude สร้าง landing page HTML + CSS ให้",
     "ใช้สีและ content จาก brand guidelines",
     "Deploy ขึ้น internet — ได้ URL *.vercel.app"],
    "Landing page live บน internet ทดสอบบนมือถือได้",
    "08", "Publisher + Wrap-up",
    ["สร้าง content-publisher skill (format per platform)",
     "สร้าง /content-engine master command",
     "รัน full pipeline ครั้งแรก",
     "เรียนรู้ 3 แนวทางปรับใช้ต่อ"],
    "/content-engine command พร้อมใช้งาน",
    "Lessons 7–8"
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Section Header: Concepts
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, VIOLET, None)
add_text(s, "ส่วนที่ 03", LM, Inches(2.0), CW, Inches(0.45), 14, True, VIOLET_L)
add_text(s, "แนวคิดหลัก\nที่คุณจะเข้าใจ", LM, Inches(2.55), CW, Inches(1.8), 40, True, TEXT)
tagline_bar(s, Inches(4.55))
add_text(s, "ไม่ใช่แค่ใช้เครื่องมือ — แต่เข้าใจว่ามันทำงานอย่างไร",
         LM, Inches(4.75), CW, Inches(0.5), 16, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Core Pattern
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Pattern ที่ใช้ซ้ำได้ทุกระบบ")
tagline_bar(s)

flow_items = [
    ("Brand Context", "CLAUDE.md + guidelines\nทุก agent อ่านจากที่เดียวกัน", VIOLET),
    ("Skills",        "SKILL.md ต่อหนึ่งงาน\nบอกว่าต้องการงานแบบไหน", CYAN),
    ("MCPs",          "เชื่อมกับระบบภายนอก\nNocoDB, HeyGen, Vercel", GREEN),
    ("Slash Commands","กด / ครั้งเดียว\nรัน workflow ทั้งหมด", AMBER),
]
n = len(flow_items)
box_w = Inches(2.9)
arrow_w = Inches(0.35)
total_needed = box_w * n + arrow_w * (n - 1)
start_x = (SLIDE_W - total_needed) / 2
y_flow = Inches(2.2)
fh = Inches(2.2)

x = start_x
for i, (title, desc, col) in enumerate(flow_items):
    card(s, x, y_flow, box_w, fh, col)
    add_text(s, title, x + Inches(0.15), y_flow + Inches(0.2),
             box_w - Inches(0.2), Inches(0.5), 16, True, TEXT, PP_ALIGN.CENTER)
    add_text(s, desc, x + Inches(0.12), y_flow + Inches(0.8),
             box_w - Inches(0.2), Inches(1.2), 12, False, MUTED, PP_ALIGN.CENTER)
    x += box_w
    if i < n - 1:
        add_text(s, "→", x, y_flow + Inches(0.75), arrow_w,
                 Inches(0.7), 22, True, VIOLET, PP_ALIGN.CENTER)
        x += arrow_w

add_rect(s, LM, Inches(4.7), CW, Inches(0.55), RGBColor(0x12, 0x0A, 0x28), VIOLET)
add_text(s, "เรียนรู้ pattern นี้ครั้งเดียว — สร้างระบบ AI ได้ไม่จำกัด",
         LM + Inches(0.2), Inches(4.78), CW - Inches(0.3), Inches(0.38),
         13, False, VIOLET_L, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CLAUDE.md
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "CLAUDE.md")
tagline_bar(s)
subheading(s, "บริบทถาวรที่ Claude อ่านทุก session อัตโนมัติ")

# Code block mock
cx = LM
card(s, cx, Inches(1.85), Inches(5.9), Inches(3.2), BORDER)
code_lines = [
    ("# marketing-agency/CLAUDE.md", MUTED),
    ("", TEXT),
    ("## บริบทธุรกิจ", VIOLET_L),
    ("แบรนด์: [ชื่อแบรนด์คุณ]", CODE),
    ("สิ่งที่ขาย: [สินค้า/บริการ]", CODE),
    ("กลุ่มเป้าหมาย: [ลูกค้าคุณ]", CODE),
    ("", TEXT),
    ("## ไฟล์สำคัญ", VIOLET_L),
    ("brand/brand-guidelines.md", CODE),
    (".claude/skills/", CODE),
]
for j, (line, col) in enumerate(code_lines):
    add_text(s, line, cx + Inches(0.2),
             Inches(2.0) + Inches(0.28) * j,
             Inches(5.5), Inches(0.3), 11, False, col)

# Right col
rx = LM + Inches(6.2)
for j, (title, desc, col) in enumerate([
    ("ไม่ต้องอธิบายโปรเจกต์ใหม่ทุกครั้ง",
     "Claude จำบริบทธุรกิจของคุณข้ามทุก session", VIOLET),
    ("ทุก agent อ่าน briefing เดียวกัน",
     "researcher, writer, designer — รู้จักแบรนด์คุณเท่ากัน", VIOLET),
    ("เหมือน SOP ที่พนักงานใหม่ต้องอ่าน",
     "เขียนครั้งเดียว ใช้ได้ตลอด แก้ไขได้ตลอด", VIOLET),
]):
    y = Inches(1.85) + Inches(1.12) * j
    card(s, rx, y, Inches(5.9), Inches(1.0), col)
    add_text(s, title, rx + Inches(0.2), y + Inches(0.1),
             Inches(5.5), Inches(0.4), 13, True, TEXT)
    add_text(s, desc, rx + Inches(0.2), y + Inches(0.52),
             Inches(5.5), Inches(0.38), 11, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Skills
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Skills")
tagline_bar(s)
subheading(s, "สอน Claude ว่าคุณต้องการงานแบบไหน")

# Analogy callout
add_rect(s, LM, Inches(1.85), CW, Inches(0.55), RGBColor(0x06, 0x1A, 0x20), CYAN)
add_text(s, "เหมือน SOP ของพนักงาน — เขียนครั้งเดียว ใช้ซ้ำได้เสมอ พร้อม context ของแบรนด์คุณ",
         LM + Inches(0.2), Inches(1.93), CW - Inches(0.3), Inches(0.38),
         13, False, CYAN, PP_ALIGN.CENTER)

skills = [
    (".claude/skills/content-researcher/SKILL.md",
     "รู้ว่าต้องหาข้อมูลจากไหน ออกแบบไว้สำหรับ niche ของคุณ"),
    (".claude/skills/content-writer/SKILL.md",
     "รู้ว่าต้องเขียนสไตล์ไหน เสียงแบรนด์ไหน"),
    (".claude/skills/video-maker/SKILL.md",
     "รู้ว่าต้องใช้ avatar ไหน voice ไหน format ไหน"),
    (".claude/skills/content-designer/SKILL.md",
     "รู้สี typography และ style direction ของแบรนด์คุณ"),
]
for i, (skill_path, desc) in enumerate(skills):
    y = Inches(2.65) + Inches(1.0) * i
    card(s, LM, y, CW, Inches(0.88), BORDER)
    add_text(s, skill_path, LM + Inches(0.2), y + Inches(0.1),
             Inches(7.5), Inches(0.38), 12, False, CODE)
    add_text(s, "→  " + desc, LM + Inches(0.2), y + Inches(0.5),
             CW - Inches(0.3), Inches(0.32), 12, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — MCPs
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "MCPs — Model Context Protocol")
tagline_bar(s)
subheading(s, "ให้ Claude เชื่อมกับระบบภายนอก ทำงานกับโลกภายนอกได้")

# Center Claude box
center_x = (SLIDE_W - Inches(2.2)) / 2
center_y = Inches(2.5)
add_rect(s, center_x, center_y, Inches(2.2), Inches(2.2), VIOLET, None)
add_text(s, "Claude\nCode", center_x, center_y + Inches(0.6),
         Inches(2.2), Inches(0.9), 20, True, WHITE, PP_ALIGN.CENTER)

# Services
services = [
    ("NocoDB",  "เพิ่ม/อ่าน/แก้ไข rows ใน database", GREEN,  LM,                        Inches(3.1)),
    ("HeyGen",  "สร้างวิดีโอ AI avatar",              VIOLET, LM,                        Inches(4.3)),
    ("Vercel",  "Deploy เว็บขึ้น internet",           CYAN,   LM,                        Inches(5.5)),
]
for name, desc, col, sx, sy in services:
    card(s, sx, sy, Inches(3.5), Inches(0.88), col)
    add_text(s, name, sx + Inches(0.15), sy + Inches(0.08),
             Inches(3.2), Inches(0.38), 14, True, TEXT)
    add_text(s, desc, sx + Inches(0.15), sy + Inches(0.48),
             Inches(3.2), Inches(0.32), 11, False, MUTED)
    # Arrow right
    add_text(s, "⇄", sx + Inches(3.6), sy + Inches(0.28),
             Inches(0.6), Inches(0.4), 18, True, col, PP_ALIGN.CENTER)

add_rect(s, LM, Inches(6.55), CW, Inches(0.52), RGBColor(0x0A, 0x10, 0x18), BORDER)
add_text(s, "ไม่มี MCP = Claude รู้แค่ไฟล์ในเครื่อง   |   มี MCP = Claude ทำงานกับโลกภายนอกได้",
         LM + Inches(0.2), Inches(6.63), CW - Inches(0.3), Inches(0.35),
         12, False, MUTED, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Slash Commands
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Slash Commands")
tagline_bar(s)
subheading(s, "พิมพ์ / ครั้งเดียว — รัน workflow ที่ซับซ้อนทั้งหมด")

commands = [
    ("/brand",          "ถามและอัปเดต brand guidelines",
     "เมื่อแบรนด์คุณเปลี่ยน Claude จะถามคำถามและอัปเดตไฟล์"),
    ("/research",       "วิจัยหัวข้อสำหรับ content pillar ที่เลือก",
     "หา trending topics, competitor angles, audience questions"),
    ("/content-engine", "รัน pipeline ทั้งหมด",
     "research → calendar → write → design → format — ครบในคำสั่งเดียว"),
]
for i, (cmd, title, desc) in enumerate(commands):
    y = Inches(1.85) + Inches(1.55) * i
    card(s, LM, y, CW, Inches(1.38), BORDER)
    # Command pill
    add_rect(s, LM + Inches(0.15), y + Inches(0.2), Inches(2.8), Inches(0.55),
             SURFACE2, VIOLET)
    add_text(s, cmd, LM + Inches(0.25), y + Inches(0.25),
             Inches(2.6), Inches(0.42), 15, True, CODE)
    add_text(s, title, LM + Inches(3.15), y + Inches(0.2),
             CW - Inches(3.25), Inches(0.4), 14, True, TEXT)
    add_text(s, desc, LM + Inches(3.15), y + Inches(0.65),
             CW - Inches(3.25), Inches(0.55), 12, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Section Header: Adaptations
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, VIOLET, None)
add_text(s, "ส่วนที่ 04", LM, Inches(2.0), CW, Inches(0.45), 14, True, VIOLET_L)
add_text(s, "3 สิ่งที่คุณ\nสามารถทำต่อได้", LM, Inches(2.55), CW, Inches(1.8), 40, True, TEXT)
tagline_bar(s, Inches(4.55))
add_text(s, "บทเรียนที่ 8 สอนการปรับระบบนี้สำหรับงานจริง",
         LM, Inches(4.75), CW, Inches(0.5), 16, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Adaptation 1
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, LM, TM, Inches(2.0), Inches(0.38), SURFACE, VIOLET)
add_text(s, "Adaptation 01", LM + Inches(0.1), TM + Inches(0.08),
         Inches(1.8), Inches(0.25), 10, True, VIOLET_L)
heading(s, "ระบบ Onboarding ลูกค้า", Inches(0.95))
tagline_bar(s, Inches(1.55))
subheading(s, "สำหรับ Freelancer, Agency, Consultant", Inches(1.72), 14, MUTED)

steps_adapt = [
    ("1.", "คัดลอกโฟลเดอร์ marketing-agency ทั้งหมด",
     "สร้างโฟลเดอร์ใหม่ต่อลูกค้า 1 ราย โครงสร้างทุกอย่างเหมือนกัน"),
    ("2.", "เปลี่ยนแค่ brand-guidelines.md",
     "ใส่ข้อมูลแบรนด์ลูกค้า — ทุก agent จะทำงานตาม brief ใหม่ทันที"),
    ("3.", "รัน /content-engine",
     "ได้ content calendar 1 เดือนสำหรับลูกค้าราย 1 คน พร้อม draft ทันที"),
]
for i, (num, title, desc) in enumerate(steps_adapt):
    y = Inches(2.2) + Inches(1.35) * i
    card(s, LM, y, CW, Inches(1.22), BORDER)
    add_text(s, num, LM + Inches(0.15), y + Inches(0.25),
             Inches(0.55), Inches(0.55), 26, True, VIOLET_L)
    add_text(s, title, LM + Inches(0.8), y + Inches(0.12),
             CW - Inches(0.95), Inches(0.42), 15, True, TEXT)
    add_text(s, desc, LM + Inches(0.8), y + Inches(0.6),
             CW - Inches(0.95), Inches(0.5), 12, False, MUTED)

add_rect(s, LM, Inches(6.3), CW, Inches(0.52), RGBColor(0x06, 0x20, 0x17), GREEN)
add_text(s, "สร้างระบบครั้งเดียว — ใช้ได้ทุกลูกค้า",
         LM + Inches(0.2), Inches(6.38), CW - Inches(0.3), Inches(0.35),
         13, True, GREEN, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Adaptations 2 & 3
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Adaptations 2 & 3")
tagline_bar(s)

adaptations = [
    ("02", "Product Launch Campaign",
     ["ปรับ content-calendar skill ให้รับ launch date",
      "Generate teaser → announcement → follow-up sequence",
      "Writer สร้างเนื้อหา launch-specific ทุกชิ้นอัตโนมัติ",
      "Designer สร้างรูปธีม launch สอดคล้องกัน"]),
    ("03", "Seasonal Campaign",
     ["วิจัย seasonal trends ล่วงหน้าหลายเดือน",
      "Calendar เต็มไปด้วย timely content hooks",
      "Writer และ Designer สร้าง content เป็น batch",
      "คอนเทนต์พร้อมก่อน season จริง"]),
]
for xi, (num, title, items) in enumerate(adaptations):
    x = LM if xi == 0 else LM + Inches(6.2)
    card(s, x, Inches(1.7), Inches(5.9), Inches(5.1), BORDER)
    add_text(s, num, x + Inches(0.15), Inches(1.85),
             Inches(1.2), Inches(0.65), 32, True, VIOLET_L)
    add_text(s, title, x + Inches(0.15), Inches(2.6),
             Inches(5.6), Inches(0.45), 15, True, TEXT)
    add_rect(s, x + Inches(0.1), Inches(3.12), Inches(5.7), Inches(0.02), BORDER, None)
    for j, item in enumerate(items):
        add_text(s, "—  " + item, x + Inches(0.15),
                 Inches(3.22) + Inches(0.72) * j,
                 Inches(5.6), Inches(0.62), 12, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Pattern repeatable
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Pattern เดิม  ใช้ได้ทุกโปรเจกต์")
tagline_bar(s)

flow_labels = [
    ("CLAUDE.md", "บริบทถาวร"),
    ("Skills",    "วิธีทำงาน"),
    ("MCPs",      "ระบบนอก"),
    ("Commands",  "Trigger"),
]
n = len(flow_labels)
box_w = Inches(2.8)
arrow_w2 = Inches(0.3)
start_x2 = LM
y2 = Inches(2.4)
fh2 = Inches(1.8)

x = start_x2
for i, (lbl, sub) in enumerate(flow_labels):
    card(s, x, y2, box_w, fh2, VIOLET)
    add_text(s, lbl, x, y2 + Inches(0.38), box_w, Inches(0.5),
             18, True, TEXT, PP_ALIGN.CENTER)
    add_text(s, sub, x, y2 + Inches(0.95), box_w, Inches(0.45),
             12, False, MUTED, PP_ALIGN.CENTER)
    x += box_w
    if i < n - 1:
        add_text(s, "→", x, y2 + Inches(0.6), arrow_w2,
                 Inches(0.55), 20, True, VIOLET_L, PP_ALIGN.CENTER)
        x += arrow_w2

add_text(s, "คุณสร้างระบบนี้แล้ว 1 ระบบ ระบบถัดไปจะใช้เวลาครึ่งหนึ่ง เพราะรู้ pattern แล้ว",
         LM, Inches(4.5), CW, Inches(0.5), 14, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Section Header: Limits (amber)
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, AMBER, None)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x12, 0x0C, 0x00), None)
fill_bg(s, RGBColor(0x12, 0x0C, 0x00))
add_text(s, "ส่วนที่ 05", LM, Inches(2.0), CW, Inches(0.45), 14, True, AMBER)
add_text(s, "สิ่งที่ระบบนี้\nยังทำไม่ได้", LM, Inches(2.55), CW, Inches(1.8), 40, True, TEXT)
tagline_bar(s, Inches(4.55), AMBER)
add_text(s, "พูดตรงๆ เพื่อให้คุณรู้ว่ากำลังลงทุนอะไร",
         LM, Inches(4.75), CW, Inches(0.5), 16, False, MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Limits table
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "ความจริงที่ควรรู้ก่อนเริ่ม")
tagline_bar(s)

can_items = [
    "วิจัย, วางแผน, เขียน, ออกแบบ, ทำวิดีโอ",
    "สร้าง landing page และ deploy",
    "Format content สำหรับแต่ละ platform",
    "สร้าง content calendar ล่วงหน้า 1 เดือน",
    "รัน pipeline ทั้งหมดด้วยคำสั่งเดียว",
]
cant_items = [
    "ตรวจและอนุมัติก่อนโพสต์",
    "โพสต์จริงลงแพลตฟอร์ม (copy-paste)",
    "เชื่อมต่อโดยตรงกับ social media API",
    "วิเคราะห์ performance หลังโพสต์",
    "A/B testing และ analytics dashboard",
]

col_w3 = Inches(5.9)
# Header row
add_rect(s, LM, Inches(1.7), col_w3, Inches(0.5),
         RGBColor(0x06, 0x20, 0x17), GREEN)
add_text(s, "สิ่งที่ระบบทำได้",
         LM + Inches(0.15), Inches(1.77), col_w3 - Inches(0.2), Inches(0.38),
         14, True, GREEN)

rx3 = LM + Inches(6.2)
add_rect(s, rx3, Inches(1.7), col_w3, Inches(0.5),
         RGBColor(0x1A, 0x10, 0x00), AMBER)
add_text(s, "สิ่งที่ยังต้องทำเอง",
         rx3 + Inches(0.15), Inches(1.77), col_w3 - Inches(0.2), Inches(0.38),
         14, True, AMBER)

for i, (can, cant) in enumerate(zip(can_items, cant_items)):
    y_row = Inches(2.3) + Inches(0.78) * i
    bg_col = SURFACE if i % 2 == 0 else SURFACE2
    add_rect(s, LM, y_row, col_w3, Inches(0.72), bg_col, BORDER)
    add_text(s, "✓  " + can, LM + Inches(0.12), y_row + Inches(0.14),
             col_w3 - Inches(0.2), Inches(0.5), 12, False, GREEN)
    add_rect(s, rx3, y_row, col_w3, Inches(0.72), bg_col, BORDER)
    add_text(s, "✗  " + cant, rx3 + Inches(0.12), y_row + Inches(0.14),
             col_w3 - Inches(0.2), Inches(0.5), 12, False, AMBER)

add_rect(s, LM, Inches(6.3), CW, Inches(0.5), RGBColor(0x1A, 0x10, 0x00), AMBER)
add_text(s, "ระบบนี้จัดการ production work — คุณจัดการ judgement calls",
         LM + Inches(0.2), Inches(6.38), CW - Inches(0.3), Inches(0.35),
         13, False, AMBER, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Who it's for
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "คอร์สนี้เหมาะกับคุณถ้า...")
tagline_bar(s)

fit = [
    "เจ้าของ SME ที่ทำการตลาดเอง",
    "Solopreneur ที่อยากได้ output มากขึ้น",
    "Marketing manager ที่ต้องการขยายผล",
    "อยากเข้าใจ Claude Code จากการสร้างจริง",
    "มีแบรนด์หรือธุรกิจที่ต้องการ content",
]
notfit = [
    "คาดหวัง fully automated ทันที",
    "ต้องการ multi-user team access",
    "ต้องการ analytics dashboard พร้อมใช้",
    "ไม่มี Claude Code ติดตั้ง",
    "ไม่มีแบรนด์หรือธุรกิจในใจ",
]

col_w4 = Inches(5.9)
# Green zone
add_rect(s, LM, Inches(1.7), col_w4, Inches(4.9),
         RGBColor(0x06, 0x20, 0x17), GREEN)
add_text(s, "เหมาะมาก", LM + Inches(0.2), Inches(1.85),
         col_w4 - Inches(0.3), Inches(0.42), 14, True, GREEN)
for i, t in enumerate(fit):
    add_text(s, "✓  " + t, LM + Inches(0.2), Inches(2.38) + Inches(0.75) * i,
             col_w4 - Inches(0.3), Inches(0.62), 13, False, TEXT)

# Amber zone
rx4 = LM + Inches(6.2)
add_rect(s, rx4, Inches(1.7), col_w4, Inches(4.9),
         RGBColor(0x1A, 0x10, 0x00), AMBER)
add_text(s, "อาจไม่เหมาะ", rx4 + Inches(0.2), Inches(1.85),
         col_w4 - Inches(0.3), Inches(0.42), 14, True, AMBER)
for i, t in enumerate(notfit):
    add_text(s, "✗  " + t, rx4 + Inches(0.2), Inches(2.38) + Inches(0.75) * i,
             col_w4 - Inches(0.3), Inches(0.62), 13, False, TEXT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — What to prepare
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "สิ่งที่ต้องเตรียม")
tagline_bar(s)

col_w5 = Inches(3.8)
cols = [
    ("ต้องมีก่อนเริ่ม", GREEN,
     ["Claude Code ติดตั้งในเครื่อง",
      "แบรนด์หรือธุรกิจที่อยากทำ marketing",
      "เวลา ~2.5 ชั่วโมง (ทำในวันเดียว)"]),
    ("ตั้งค่าระหว่างเรียน", CYAN,
     ["OpenRouter account (Lesson 1)",
      "NocoDB account (Lesson 3)",
      "HeyGen account (Lesson 6)",
      "Vercel account (Lesson 7)"]),
    ("ไม่ต้องมี", RGBColor(0xFC, 0xA5, 0xA5),
     ["ประสบการณ์เขียนโค้ด",
      "ทักษะดีไซน์หรือทำวิดีโอ",
      "ความรู้เรื่อง git, npm, terminal",
      "ทีมหรือพนักงาน"]),
]
border_cols = [GREEN, CYAN, RED]
for xi, ((header, hcol, items), bcol) in enumerate(zip(cols, border_cols)):
    x5 = LM + (col_w5 + Inches(0.2)) * xi
    card(s, x5, Inches(1.7), col_w5, Inches(5.0), bcol)
    add_text(s, header, x5 + Inches(0.15), Inches(1.85),
             col_w5 - Inches(0.2), Inches(0.42), 13, True, hcol)
    add_rect(s, x5 + Inches(0.1), Inches(2.32), col_w5 - Inches(0.2), Inches(0.02), BORDER, None)
    prefix = "✗" if xi == 2 else "•"
    prefix_col = RED if xi == 2 else CYAN
    for j, item in enumerate(items):
        add_text(s, prefix + "  " + item,
                 x5 + Inches(0.15), Inches(2.45) + Inches(0.85) * j,
                 col_w5 - Inches(0.2), Inches(0.75), 12, False,
                 TEXT if xi < 2 else MUTED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 31 — How to start
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "วิธีเริ่มเรียน")
tagline_bar(s)

start_steps = [
    ("1.", "เปิดโฟลเดอร์นี้ใน Claude Code",
     "ไปที่ File → Open Folder → เลือกโฟลเดอร์ marketing-agency-tutorial"),
    ("2.", "พิมพ์คำสั่งด้านล่างในช่อง chat",
     "Claude จะอ่าน lesson file และกลายเป็นครูผู้สอนทันที"),
    ("3.", "ตอบคำถามและทำตามที่ Claude แนะนำ",
     "คุณไม่ได้อ่านเอกสาร — คุณกำลังสร้างระบบจริงสำหรับธุรกิจของคุณ"),
]
for i, (num, title, desc) in enumerate(start_steps):
    y = Inches(1.72) + Inches(1.22) * i
    card(s, LM, y, CW, Inches(1.1), BORDER)
    add_text(s, num, LM + Inches(0.15), y + Inches(0.22),
             Inches(0.55), Inches(0.52), 24, True, VIOLET_L)
    add_text(s, title, LM + Inches(0.8), y + Inches(0.1),
             CW - Inches(0.95), Inches(0.4), 15, True, TEXT)
    add_text(s, desc, LM + Inches(0.8), y + Inches(0.56),
             CW - Inches(0.95), Inches(0.42), 12, False, MUTED)

# Big command pill
cmd_w = Inches(5.0)
cmd_x = (SLIDE_W - cmd_w) / 2
add_rect(s, cmd_x, Inches(5.42), cmd_w, Inches(0.72),
         RGBColor(0x1A, 0x0D, 0x35), VIOLET)
add_text(s, "Start Lesson 1",
         cmd_x, Inches(5.5), cmd_w, Inches(0.55),
         22, True, CODE, PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 32 — Closing CTA
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
# Subtle glow bg
add_rect(s, Inches(6.5), Inches(-1.5), Inches(8), Inches(8),
         RGBColor(0x12, 0x06, 0x24), None)

add_text(s, "Claude Code Course",
         LM, TM, Inches(4), Inches(0.38), 11, False, VIOLET_L)

add_text(s, "Private Marketing\nAgency Blueprint",
         LM, Inches(1.0), Inches(9), Inches(2.1), 46, True, TEXT)

tagline_bar(s, Inches(3.3))

add_text(s, "8 บทเรียน  ·  ~2.5 ชั่วโมง  ·  ผลลัพธ์จริงสำหรับธุรกิจของคุณ",
         LM, Inches(3.55), CW, Inches(0.45), 14, False, MUTED)

closing_pts = [
    ("✓", "ระบบ AI marketing ที่ทำงานได้จริงสำหรับแบรนด์คุณ"),
    ("✓", "ไม่ต้องเขียนโค้ด ไม่ต้องจ้างทีม"),
    ("✓", "Pattern ที่นำไปสร้างระบบต่อไปได้เอง"),
]
for i, (sym, text) in enumerate(closing_pts):
    add_text(s, sym, LM, Inches(4.25) + Inches(0.58) * i,
             Inches(0.4), Inches(0.45), 16, True, GREEN)
    add_text(s, text, LM + Inches(0.45), Inches(4.25) + Inches(0.58) * i,
             CW - Inches(0.5), Inches(0.45), 15, False, TEXT)

# CTA box
add_rect(s, LM, Inches(6.38), CW, Inches(0.68), VIOLET, None)
add_text(s, "เปิดโฟลเดอร์ใน Claude Code แล้วพิมพ์  \"Start Lesson 1\"",
         LM, Inches(6.48), CW, Inches(0.5), 15, True, WHITE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out_path = "/Users/watchara/Documents/me/skill-private-marketing-agency-tutorial/presentation/slides.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
